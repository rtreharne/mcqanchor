import hashlib
import json
import re
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from django.conf import settings
from django.db import transaction
from markdown import markdown
from openai import OpenAI
from openpyxl import load_workbook
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from standalone.models import ContentAsset, ContentChunk, LearningObjective


SUPPORTED_EXTENSIONS = {
    ".html",
    ".docx",
    ".pdf",
    ".txt",
    ".r",
    ".py",
    ".ipynb",
    ".rmd",
    ".md",
    ".pptx",
    ".xlsx",
}


def normalize_text(value: str) -> str:
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _plain_text_from_markdown(text: str) -> str:
    return BeautifulSoup(markdown(text), "html.parser").get_text("\n")


def extract_text_from_asset(asset: ContentAsset) -> str:
    ext = asset.extension.lower()
    file_path = Path(asset.file.path)

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported content type: {ext}")

    if ext in {".txt", ".r", ".py"}:
        return normalize_text(file_path.read_text(encoding="utf-8", errors="ignore"))

    if ext in {".md", ".rmd"}:
        return normalize_text(_plain_text_from_markdown(file_path.read_text(encoding="utf-8", errors="ignore")))

    if ext == ".html":
        return normalize_text(BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "html.parser").get_text("\n"))

    if ext == ".ipynb":
        notebook = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
        cells = []
        for cell in notebook.get("cells", []):
            source = "".join(cell.get("source", []))
            if cell.get("cell_type") == "markdown":
                cells.append(_plain_text_from_markdown(source))
            else:
                cells.append(source)
        return normalize_text("\n\n".join(cells))

    if ext == ".docx":
        document = Document(file_path)
        return normalize_text("\n".join(paragraph.text for paragraph in document.paragraphs))

    if ext == ".pdf":
        reader = PdfReader(str(file_path))
        return normalize_text("\n".join(page.extract_text() or "" for page in reader.pages))

    if ext == ".pptx":
        deck = Presentation(str(file_path))
        slides = []
        for slide in deck.slides:
            slide_lines = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            if slide_lines:
                slides.append("\n".join(slide_lines))
        return normalize_text("\n\n".join(slides))

    if ext == ".xlsx":
        workbook = load_workbook(str(file_path), data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell not in (None, "")]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"{sheet.title}\n" + "\n".join(rows))
        return normalize_text("\n\n".join(sheets))

    raise ValueError(f"Unsupported content type: {ext}")


def _split_long_segment(text: str, target_size: int) -> list[str]:
    if len(text) <= target_size:
        return [text]

    sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", text) if segment.strip()]
    if len(sentences) > 1:
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0
        for sentence in sentences:
            if len(sentence) > target_size:
                if current:
                    chunks.append(" ".join(current))
                    current = []
                    current_len = 0
                chunks.extend(_split_long_segment(sentence, target_size))
                continue
            projected_len = current_len + len(sentence) + (1 if current else 0)
            if current and projected_len > target_size:
                chunks.append(" ".join(current))
                current = [sentence]
                current_len = len(sentence)
            else:
                current.append(sentence)
                current_len = projected_len
        if current:
            chunks.append(" ".join(current))
        return chunks

    words = text.split()
    if not words:
        return []

    chunks = []
    current_words: list[str] = []
    current_len = 0
    for word in words:
        if len(word) > target_size:
            if current_words:
                chunks.append(" ".join(current_words))
                current_words = []
                current_len = 0
            chunks.extend([word[index : index + target_size] for index in range(0, len(word), target_size)])
            continue
        projected_len = current_len + len(word) + (1 if current_words else 0)
        if current_words and projected_len > target_size:
            chunks.append(" ".join(current_words))
            current_words = [word]
            current_len = len(word)
        else:
            current_words.append(word)
            current_len = projected_len
    if current_words:
        chunks.append(" ".join(current_words))
    return chunks


def chunk_text(text: str, target_size: int = 1200) -> list[str]:
    paragraphs = [segment.strip() for segment in text.split("\n\n") if segment.strip()]
    if not paragraphs:
        return []

    chunks = []
    current = []
    current_len = 0
    for paragraph in paragraphs:
        paragraph_segments = _split_long_segment(paragraph, target_size)
        for segment in paragraph_segments:
            projected_len = current_len + len(segment) + (2 if current else 0)
            if current and projected_len > target_size:
                chunks.append("\n\n".join(current))
                current = [segment]
                current_len = len(segment)
            else:
                current.append(segment)
                current_len = projected_len if current_len else len(segment)
    if current:
        chunks.append("\n\n".join(current))
    return chunks


def sanitize_learning_objective(text: str) -> str:
    value = text.strip().strip("\"'`")
    value = value.replace("\u2013", "-").replace("\u2014", "-").replace("\u2212", "-")
    value = value.replace("\u2018", "'").replace("\u2019", "'").replace("\u201c", '"').replace("\u201d", '"')
    value = re.sub(r"^\s*[•●◦▪·*-]+\s*", "", value)
    value = re.sub(r"^\s*(?:learning\s*objective|objective|outcome|lo)\s*[:#-]?\s*\d+[a-z]?\s*[:.)-]?\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*\(?\d+[a-z]?\)?\s*[.)-]\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*\(?[ivxlcdm]+\)?\s*[.)-]\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"^\s*[A-Z]\s*[.)-]\s*", "", value)
    value = value.strip(" -:\t")
    value = re.sub(r"\s+", " ", value)
    value = re.sub(r"[;:,.\-]+$", "", value)
    if not value:
        return ""
    if value[0].isalpha():
        value = value[0].upper() + value[1:]
    return value


def sanitize_summary(text: str) -> str:
    value = normalize_text(text)
    value = value.replace("\u2013", "-").replace("\u2014", "-")
    value = value.replace("\u2018", "'").replace("\u2019", "'").replace("\u201c", '"').replace("\u201d", '"')
    prefix_patterns = (
        r"^(?:the\s+)?(?:teaching\s+content|uploaded\s+teaching\s+material|uploaded\s+material|course\s+material|course\s+content|content|material)\s+(?:provides|offers|covers|introduces|explores|discusses|focuses\s+on|is\s+about)\s+",
        r"^(?:this\s+)?(?:teaching\s+content|uploaded\s+teaching\s+material|uploaded\s+material|course\s+material|course\s+content|content|material)\s+(?:provides|offers|covers|introduces|explores|discusses|focuses\s+on|is\s+about)\s+",
    )
    for pattern in prefix_patterns:
        value = re.sub(pattern, "", value, flags=re.IGNORECASE)
    value = re.sub(r"^(?:an?|the)\s+overview\s+of\s+", "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+([,.;:!?])", r"\1", value)
    value = re.sub(r"[ \t]+", " ", value)
    value = value.strip()
    if value and value[0].isalpha():
        value = value[0].upper() + value[1:]
    return value


def _dedupe_texts(items: list[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for item in items:
        key = re.sub(r"[^a-z0-9]+", " ", item.lower()).strip()
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(item)
    return deduped


def _objective_budget_for_text(text: str, minimum: int = 6, maximum: int = 12) -> int:
    sections = max(1, len(chunk_text(text, target_size=1200)))
    line_candidates = len([line for line in text.splitlines() if sanitize_learning_objective(line)])
    return max(minimum, min(maximum, max(sections * 2, line_candidates)))


def derive_learning_objectives(text: str, max_items: int = 8) -> list[str]:
    candidates = []
    for line in text.splitlines():
        stripped = sanitize_learning_objective(line)
        if 30 <= len(stripped) <= 180 and stripped not in candidates:
            candidates.append(stripped)
        if len(candidates) >= max_items:
            break

    if not candidates:
        sentences = re.split(r"(?<=[.!?])\s+", text)
        for sentence in sentences:
            stripped = sanitize_learning_objective(sentence)
            if 30 <= len(stripped) <= 160:
                candidates.append(stripped)
            if len(candidates) >= max_items:
                break

    return _dedupe_texts(candidates)[:max_items]


def derive_learning_objectives_with_coverage(text: str, max_items: int = 8) -> list[str]:
    sections = chunk_text(text, target_size=1200)
    if len(sections) <= 1:
        return derive_learning_objectives(text, max_items=max_items)

    per_section_limit = max(2, min(4, max_items // max(1, len(sections)) + 1))
    candidates: list[str] = []

    for section in sections:
        candidates.extend(derive_learning_objectives(section, max_items=per_section_limit))

    candidates.extend(derive_learning_objectives(text, max_items=max_items))
    return _dedupe_texts(candidates)[:max_items]


def _fallback_summary(text: str, max_sentences: int = 2, max_length: int = 320) -> str:
    normalized = normalize_text(text)
    sections = chunk_text(normalized, target_size=1200)
    representative_sentences: list[str] = []

    if len(sections) > 1:
        section_indexes = [0]
        if max_sentences > 1:
            section_indexes.append(len(sections) - 1)
        if max_sentences > 2 and len(sections) > 2:
            section_indexes.insert(1, len(sections) // 2)

        seen_sentences: set[str] = set()
        for index in section_indexes:
            section_sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", sections[index]) if segment.strip()]
            if not section_sentences:
                continue
            sentence = section_sentences[0]
            sentence_key = sentence.lower()
            if sentence_key in seen_sentences:
                continue
            seen_sentences.add(sentence_key)
            representative_sentences.append(sentence)
            if len(representative_sentences) >= max_sentences:
                break

    if representative_sentences:
        summary = " ".join(representative_sentences).strip()
    else:
        sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", normalized) if segment.strip()]
        summary = " ".join(sentences[:max_sentences]).strip()
    if not summary:
        summary = normalized[:max_length]
    if len(summary) > max_length:
        summary = summary[: max_length - 1].rsplit(" ", 1)[0].rstrip(",;:-") + "."
    if summary and summary[-1] not in ".!?":
        summary += "."
    return sanitize_summary(summary)


def _parse_json_object(raw_output: str) -> dict[str, Any]:
    if not raw_output:
        raise ValueError("OpenAI returned an empty payload.")
    try:
        return json.loads(raw_output)
    except json.JSONDecodeError:
        pass

    fenced_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw_output, re.DOTALL)
    if fenced_match:
        return json.loads(fenced_match.group(1))

    object_match = re.search(r"\{.*\}", raw_output, re.DOTALL)
    if object_match:
        return json.loads(object_match.group(0))

    raise ValueError("OpenAI did not return parseable JSON.")


def _openai_block_content_payload(text: str, max_items: int) -> dict[str, Any]:
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    prompt = f"""
Summarise the teaching content below for an educator-facing course authoring workflow.

Return strict JSON with exactly these keys:
- summary: one or two sentences in plain English
- learning_objectives: an array of 3 to {max_items} concise learning objectives

Rules:
- no numbering or bullets in the output
- write the summary as a direct course/block description, not a description of the uploaded content
- do not use phrases like "the teaching content", "the uploaded material", "the content provides", "the textbook covers", or "this material discusses"
- no ambiguous references such as "this", "that", "it", or "etc."
- each learning objective should start with a clear action verb
- remove odd characters and formatting noise
- keep the wording specific enough to make sense on its own

Content:
{text}
""".strip()
    response = client.responses.create(
        model=settings.OPENAI_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": "Return only valid JSON."}]},
            {"role": "user", "content": [{"type": "input_text", "text": prompt}]},
        ],
    )
    return _parse_json_object((getattr(response, "output_text", "") or "").strip())


def _sanitize_objective_candidates(items: list[Any], max_items: int) -> list[str]:
    if not isinstance(items, list):
        return []
    return _dedupe_texts(
        [
            sanitized
            for item in items
            if isinstance(item, str)
            for sanitized in [sanitize_learning_objective(item)]
            if sanitized
        ]
    )[:max_items]


def _openai_reduce_content_payload(section_summaries: list[str], objective_candidates: list[str], max_items: int) -> dict[str, Any]:
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    summary_lines = "\n".join(f"- {summary}" for summary in section_summaries if summary.strip())
    objective_lines = "\n".join(f"- {objective}" for objective in objective_candidates if objective.strip())
    prompt = f"""
Synthesize an educator-facing summary and learning objectives from the uploaded teaching material notes below.

Return strict JSON with exactly these keys:
- summary: one or two sentences in plain English
- learning_objectives: an array of 3 to {max_items} concise learning objectives

Rules:
- use only the supplied uploaded-material notes
- no numbering or bullets in the output
- write the summary as a direct course/block description, not a description of the uploaded notes
- do not use phrases like "the teaching content", "the uploaded material", "the content provides", "the textbook covers", or "this material discusses"
- no ambiguous references such as "this", "that", "it", or "etc."
- each learning objective should start with a clear action verb
- remove odd characters and formatting noise
- keep the wording specific enough to make sense on its own

Section summaries:
{summary_lines or "- None"}

Candidate learning objectives:
{objective_lines or "- None"}
""".strip()
    response = client.responses.create(
        model=settings.OPENAI_MODEL,
        input=[
            {"role": "system", "content": [{"type": "input_text", "text": "Return only valid JSON."}]},
            {"role": "user", "content": [{"type": "input_text", "text": prompt}]},
        ],
    )
    return _parse_json_object((getattr(response, "output_text", "") or "").strip())


def summarize_block_content(text: str, max_items: int = 6) -> tuple[str, list[str]]:
    fallback_summary = _fallback_summary(text)
    fallback_objectives = derive_learning_objectives_with_coverage(text, max_items=max_items)

    summary = fallback_summary
    objectives = fallback_objectives

    if settings.OPENAI_API_KEY:
        try:
            sections = chunk_text(text, target_size=5000)
            if len(sections) <= 1:
                payload = _openai_block_content_payload(text, max_items)
            else:
                per_section_limit = max(3, min(6, max_items))
                section_summaries: list[str] = []
                section_objectives: list[str] = []
                for section in sections:
                    section_payload = _openai_block_content_payload(section, per_section_limit)
                    section_summary = sanitize_summary(str(section_payload.get("summary", "")))
                    if section_summary:
                        section_summaries.append(section_summary)
                    section_objectives.extend(
                        _sanitize_objective_candidates(section_payload.get("learning_objectives", []), per_section_limit)
                    )
                payload = _openai_reduce_content_payload(
                    section_summaries,
                    _dedupe_texts(section_objectives),
                    max_items,
                )

            summary = sanitize_summary(str(payload.get("summary", ""))) or fallback_summary
            objectives = _sanitize_objective_candidates(payload.get("learning_objectives", []), max_items) or fallback_objectives
        except Exception:  # noqa: BLE001
            summary = fallback_summary
            objectives = fallback_objectives

    return summary, objectives


def _asset_text_for_regeneration(asset: ContentAsset) -> str:
    return asset.extracted_text or extract_text_from_asset(asset)


def resequence_learning_objectives(block, objectives: list[LearningObjective] | None = None) -> list[LearningObjective]:
    objectives = objectives or list(block.learning_objectives.order_by("position", "pk"))
    for index, objective in enumerate(objectives, start=1):
        objective.position = index
        objective.code = f"{block.order}.{index}"
    if objectives:
        LearningObjective.objects.bulk_update(objectives, ["position", "code"])
    return objectives


@transaction.atomic
def delete_block_and_resequence(block) -> None:
    course = block.course

    for asset in block.assets.all():
        asset.file.delete(save=False)

    block.delete()

    remaining_blocks = list(course.blocks.order_by("order", "created_at", "pk"))
    for index, remaining_block in enumerate(remaining_blocks, start=1):
        remaining_block.order = index
    if remaining_blocks:
        type(remaining_blocks[0]).objects.bulk_update(remaining_blocks, ["order"])

    for remaining_block in remaining_blocks:
        resequence_learning_objectives(remaining_block)

    course_fragments = [remaining_block.summary for remaining_block in remaining_blocks if remaining_block.summary.strip()]
    if course_fragments:
        course_summary, _ = summarize_block_content("\n\n".join(course_fragments), max_items=4)
        course.summary = course_summary
    else:
        course.summary = ""
    course.save(update_fields=["summary", "updated_at"])


@transaction.atomic
def move_course_block(block, direction: str) -> bool:
    blocks = list(block.course.blocks.order_by("order", "created_at", "pk"))
    try:
        current_index = next(index for index, item in enumerate(blocks) if item.pk == block.pk)
    except StopIteration:
        return False

    if direction == "up" and current_index > 0:
        swap_index = current_index - 1
    elif direction == "down" and current_index < len(blocks) - 1:
        swap_index = current_index + 1
    else:
        return False

    blocks[current_index], blocks[swap_index] = blocks[swap_index], blocks[current_index]
    for index, reordered_block in enumerate(blocks, start=1):
        reordered_block.order = index
    type(blocks[0]).objects.bulk_update(blocks, ["order"])

    for reordered_block in blocks:
        resequence_learning_objectives(reordered_block)
    return True


def _replace_block_objectives(block, source_asset: ContentAsset, objective_texts: list[str]) -> int:
    existing = list(block.learning_objectives.order_by("position", "pk"))
    created_or_updated = 0
    for index, text in enumerate(objective_texts, start=1):
        defaults = {
            "course": block.course,
            "block": block,
            "source_asset": source_asset,
            "position": index,
            "code": f"{block.order}.{index}",
            "text": text,
        }
        if index <= len(existing):
            objective = existing[index - 1]
            changed = False
            for field, value in defaults.items():
                if getattr(objective, field) != value:
                    setattr(objective, field, value)
                    changed = True
            if changed:
                objective.save(update_fields=["course", "block", "source_asset", "position", "code", "text", "updated_at"])
            created_or_updated += 1
            continue

        LearningObjective.objects.create(**defaults)
        created_or_updated += 1

    for objective in existing[len(objective_texts) :]:
        objective.delete()

    return created_or_updated


def refresh_course_summary_from_blocks(course) -> bool:
    course_fragments = [block.summary for block in course.blocks.all() if block.summary.strip()]
    if not course_fragments:
        return False

    course_summary, _ = summarize_block_content("\n\n".join(course_fragments), max_items=4)
    course.summary = course_summary
    course.save(update_fields=["summary", "updated_at"])
    return True


def regenerate_block_descriptions_and_objectives(block, progress_callback=None) -> dict[str, int]:
    if progress_callback:
        progress_callback(15)
    assets = [asset for asset in block.assets.all() if asset.include_in_generation]
    if not assets:
        if progress_callback:
            progress_callback(100)
        return {"blocks": 0, "objectives": 0}

    if progress_callback:
        progress_callback(35)
    texts = [text for asset in assets if (text := _asset_text_for_regeneration(asset))]
    combined_text = normalize_text("\n\n".join(texts))
    if not combined_text:
        if progress_callback:
            progress_callback(100)
        return {"blocks": 0, "objectives": 0}

    objective_budget = _objective_budget_for_text(combined_text)
    if progress_callback:
        progress_callback(65)
    block_summary, objectives = summarize_block_content(combined_text, max_items=objective_budget)
    if progress_callback:
        progress_callback(85)
    block.summary = block_summary
    block.save(update_fields=["summary", "updated_at"])
    objective_count = _replace_block_objectives(block, assets[0], objectives)
    refresh_course_summary_from_blocks(block.course)
    if progress_callback:
        progress_callback(100)
    return {"blocks": 1, "objectives": objective_count}


@transaction.atomic
def move_learning_objective(objective: LearningObjective, direction: str) -> bool:
    objectives = list(objective.block.learning_objectives.order_by("position", "pk"))
    try:
        current_index = next(index for index, item in enumerate(objectives) if item.pk == objective.pk)
    except StopIteration:
        return False

    if direction == "up" and current_index > 0:
        swap_index = current_index - 1
    elif direction == "down" and current_index < len(objectives) - 1:
        swap_index = current_index + 1
    else:
        return False

    objectives[current_index], objectives[swap_index] = objectives[swap_index], objectives[current_index]
    resequence_learning_objectives(objective.block, objectives)
    return True


@transaction.atomic
def delete_learning_objective_and_resequence(objective: LearningObjective) -> None:
    block = objective.block
    objective.delete()
    resequence_learning_objectives(block)


def regenerate_course_descriptions_and_objectives(course) -> dict[str, int]:
    block_count = 0
    objective_count = 0

    for block in course.blocks.prefetch_related("assets", "learning_objectives").all():
        refreshed = regenerate_block_descriptions_and_objectives(block)
        block_count += refreshed["blocks"]
        objective_count += refreshed["objectives"]

    return {"blocks": block_count, "objectives": objective_count}


def generate_embeddings(texts: list[str]) -> list[list[float]]:
    if not settings.OPENAI_API_KEY:
        return [[] for _ in texts]
    try:
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.embeddings.create(model=settings.OPENAI_EMBEDDING_MODEL, input=texts)
        return [item.embedding for item in response.data]
    except Exception:  # noqa: BLE001
        return [[] for _ in texts]


def ingest_content_asset(asset: ContentAsset) -> None:
    extracted_text = extract_text_from_asset(asset)
    asset.extracted_text = extracted_text
    asset.processing_status = ContentAsset.ProcessingStatus.PROCESSED
    asset.processing_error = ""
    asset.save(update_fields=["extracted_text", "processing_status", "processing_error", "updated_at"])

    asset.chunks.all().delete()
    asset.learning_objectives.all().delete()

    chunks = chunk_text(extracted_text)
    embeddings = generate_embeddings(chunks) if chunks else []
    for index, chunk in enumerate(chunks, start=1):
        ContentChunk.objects.create(
            asset=asset,
            course=asset.block.course,
            block=asset.block,
            ordinal=index,
            text=chunk,
            token_count=max(1, len(chunk.split())),
            embedding_model=settings.OPENAI_EMBEDDING_MODEL if embeddings and embeddings[index - 1] else "",
            embedding_vector=embeddings[index - 1] if embeddings else [],
            checksum=hashlib.sha256(chunk.encode("utf-8")).hexdigest(),
        )

    objective_budget = _objective_budget_for_text(extracted_text, minimum=4, maximum=8)
    _, objectives = summarize_block_content(extracted_text, max_items=objective_budget)
    objectives = [text for text in objectives if text]
    for index, objective in enumerate(objectives, start=1):
        LearningObjective.objects.create(
            course=asset.block.course,
            block=asset.block,
            source_asset=asset,
            position=index,
            code=f"{asset.block.order}.{index}",
            text=objective,
        )
